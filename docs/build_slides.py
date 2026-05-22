"""
Generates two PowerPoint slides matching the SKKU Capstone-Design 2026
mid-presentation template style:
  • Nearby Request feature       (with placeholder area for screenshots)
  • Project Architecture         (high-level — layers + cloud services)

Run:
    python build_slides.py

Outputs `helpvrywhere_two_slides.pptx` next to this script.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ── Template palette (sampled from the PDF) ───────────────────────────────
GREEN = RGBColor(0x1F, 0x9D, 0x55)        # CAPSTONE green
TITLE_GRAY = RGBColor(0x3A, 0x3A, 0x3A)   # title text
BODY_GRAY = RGBColor(0x55, 0x55, 0x55)    # body text
HEADER_BG = RGBColor(0xF2, 0xF2, 0xF2)    # title-bar fill
LIGHT_BORDER = RGBColor(0xC8, 0xCD, 0xD3) # box outlines
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xE2)
SOFT_BLUE_BG = RGBColor(0xE8, 0xEE, 0xF8)
SOFT_GREEN_BG = RGBColor(0xE6, 0xF7, 0xEE)
SOFT_AMBER_BG = RGBColor(0xFF, 0xF7, 0xE1)
SOFT_PINK_BG = RGBColor(0xFC, 0xE7, 0xF3)
DASH_GRAY = RGBColor(0x9E, 0xA3, 0xAA)


def set_no_line(shape):
    shape.line.fill.background()


def fill_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def add_text(shape, text, *, size=14, bold=False, italic=False, color=BODY_GRAY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def add_paragraph(tf, text, *, size=14, bold=False, italic=False,
                  color=BODY_GRAY, align=PP_ALIGN.LEFT, font="Calibri",
                  bullet=None, level=0):
    """Append a new paragraph to an existing text frame, optionally
    prefixing a bullet character ('»' for top-level, '•' for sub)."""
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    if bullet:
        run0 = p.add_run()
        run0.text = bullet + "  "
        run0.font.name = font
        run0.font.size = Pt(size)
        run0.font.bold = True
        run0.font.color.rgb = GREEN if bullet == "»" else color
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


# ── Common chrome (header bar, footer) ─────────────────────────────────────

def add_header(slide, title):
    # Header background bar — light gray spans the full slide width.
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.25),
        Inches(13.33), Inches(0.85),
    )
    fill_solid(bar, HEADER_BG)
    set_no_line(bar)

    # Green accent bar on the left of the title.
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.30), Inches(0.32),
        Inches(0.12), Inches(0.70),
    )
    fill_solid(accent, GREEN)
    set_no_line(accent)

    # Title text — serif, bold, dark gray.
    title_box = slide.shapes.add_textbox(
        Inches(0.55), Inches(0.32), Inches(12.5), Inches(0.70),
    )
    add_text(
        title_box, title,
        size=26, bold=True, color=TITLE_GRAY,
        anchor=MSO_ANCHOR.MIDDLE, font="Cambria",
    )


def add_footer(slide, page_no):
    # CAPSTONE DESIGN — bold green left.
    box = slide.shapes.add_textbox(
        Inches(0.45), Inches(7.05), Inches(4.0), Inches(0.40),
    )
    tf = add_text(box, "CAPSTONE DESIGN", size=12, bold=True, color=GREEN,
                  anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(tf, "INNOVATION  |  COLLABORATION",
                  size=8, bold=False, color=TITLE_GRAY)

    # Slide number — italic right.
    page = slide.shapes.add_textbox(
        Inches(12.7), Inches(7.10), Inches(0.50), Inches(0.30),
    )
    add_text(page, str(page_no), size=11, italic=True, color=BODY_GRAY,
             align=PP_ALIGN.RIGHT)

    # Tiny green tick on the far-right edge — matches the template.
    tick = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(13.27), Inches(7.0),
        Inches(0.06), Inches(0.45),
    )
    fill_solid(tick, GREEN)
    set_no_line(tick)


# ── Slide 1: Nearby Request feature ────────────────────────────────────────

def slide_nearby_request(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    add_header(slide, "Project implementation details –  Nearby Request")

    # ─ Left column: bullets describing the feature ──────────────────────
    intro = slide.shapes.add_textbox(
        Inches(0.45), Inches(1.30), Inches(7.6), Inches(0.55),
    )
    add_text(
        intro,
        "Real-time community map — volunteers see help requests around them",
        size=15, bold=True, color=TITLE_GRAY,
    )

    # Bulleted list — same `»` style as template.
    bullets_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(1.95), Inches(7.6), Inches(2.6),
    )
    bullets_tf = add_text(bullets_box, "", size=13, color=BODY_GRAY)
    # First paragraph — fill in via add_paragraph so we get the bullet glyph.
    bullets_tf.text = ""  # clear placeholder
    items = [
        ("»", "Live Firestore stream of all `active` requests "
                  "— no polling, instant update when a new one is posted."),
        ("»", "Volunteer's GPS + Geolocator compute distance per "
                  "request; sorted nearest first."),
        ("»", "Mapbox map preview with avatar pin (volunteer) + "
                  "numbered orange pins (requests)."),
        ("»", "Category chips filter the same stream client-side "
                  "(Groceries  ·  Transport  ·  Household  ·  Companionship)."),
        ("»", "“Full map” button → dedicated "
                  "fullscreen view with draggable bottom-sheet list."),
    ]
    # Replace the empty paragraph with the first item, then append the rest.
    p0 = bullets_tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = items[0][0] + "  "
    run0.font.name = "Calibri"; run0.font.size = Pt(13); run0.font.bold = True
    run0.font.color.rgb = GREEN
    run1 = p0.add_run()
    run1.text = items[0][1]
    run1.font.name = "Calibri"; run1.font.size = Pt(13)
    run1.font.color.rgb = BODY_GRAY
    for prefix, text in items[1:]:
        add_paragraph(bullets_tf, text, size=13, color=BODY_GRAY,
                      bullet=prefix)
    # Add some spacing between bullets.
    for p in bullets_tf.paragraphs:
        p.space_after = Pt(6)

    # ─ Flow chips below the bullets ─────────────────────────────────────
    chips_y = Inches(4.55)
    chip_w = Inches(1.45)
    chip_h = Inches(0.55)
    chip_gap = Inches(0.20)
    chips = [
        ("Firestore", SOFT_AMBER_BG),
        ("Stream", LIGHT_BORDER),
        ("+ My GPS", SOFT_BLUE_BG),
        ("Filter", LIGHT_BORDER),
        ("Map + List", SOFT_GREEN_BG),
    ]
    x = Inches(0.45)
    for i, (label, color) in enumerate(chips):
        # Chip
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, chips_y, chip_w, chip_h,
        )
        fill_solid(chip, color)
        chip.line.color.rgb = LIGHT_BORDER
        chip.line.width = Pt(0.75)
        # Label centered
        tf = chip.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Calibri"; r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = TITLE_GRAY
        # Arrow between chips (skip after last)
        if i < len(chips) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                x + chip_w + Emu(20000),
                chips_y + Inches(0.18),
                chip_gap - Emu(40000), Inches(0.20),
            )
            fill_solid(arrow, BODY_GRAY)
            set_no_line(arrow)
        x = x + chip_w + chip_gap

    # Tech callout
    note = slide.shapes.add_textbox(
        Inches(0.45), Inches(5.25), Inches(7.6), Inches(1.5),
    )
    note_tf = add_text(
        note,
        "Built with:",
        size=11, bold=True, color=TITLE_GRAY,
    )
    add_paragraph(
        note_tf,
        "Flutter  ·  Cloud Firestore  ·  Mapbox Maps SDK  "
        "·  Geolocator  ·  Firebase Auth",
        size=11, color=BODY_GRAY,
    )
    add_paragraph(note_tf, "", size=4)
    add_paragraph(
        note_tf,
        "Distance is computed client-side with Geolocator.distanceBetween, "
        "so it stays accurate even when the user moves — no extra "
        "Firestore reads required.",
        size=10, italic=True, color=BODY_GRAY,
    )

    # ─ Right column: photo placeholders ─────────────────────────────────
    # Big screenshot placeholder
    ph1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.40), Inches(1.30),
        Inches(4.55), Inches(4.10),
    )
    ph1.fill.solid()
    ph1.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    ph1.line.color.rgb = DASH_GRAY
    ph1.line.dash_style = 7  # dash
    ph1.line.width = Pt(1.25)
    add_text(
        ph1,
        "\U0001F4F1  Screenshot:\nMap preview + nearby cards",
        size=12, italic=True, color=DASH_GRAY,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Two smaller side-by-side placeholders
    ph2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.40), Inches(5.55),
        Inches(2.20), Inches(1.40),
    )
    ph2.fill.solid()
    ph2.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    ph2.line.color.rgb = DASH_GRAY
    ph2.line.dash_style = 7
    ph2.line.width = Pt(1.25)
    add_text(
        ph2, "Full map", size=11, italic=True, color=DASH_GRAY,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    ph3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(10.75), Inches(5.55),
        Inches(2.20), Inches(1.40),
    )
    ph3.fill.solid()
    ph3.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFC)
    ph3.line.color.rgb = DASH_GRAY
    ph3.line.dash_style = 7
    ph3.line.width = Pt(1.25)
    add_text(
        ph3, "Request detail", size=11, italic=True, color=DASH_GRAY,
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    add_footer(slide, 1)


# ── Brand colors (for the Nearby-Request architecture icons) ──────────────
FLUTTER_BLUE = RGBColor(0x02, 0x56, 0x9B)
FIREBASE_AMBER = RGBColor(0xF5, 0x82, 0x0C)
FIREBASE_YELLOW = RGBColor(0xFF, 0xCB, 0x2B)
MAPBOX_NAVY = RGBColor(0x0F, 0x1B, 0x30)
MAPBOX_BLUE = RGBColor(0x44, 0x6E, 0xCC)
GPS_GREEN = RGBColor(0x34, 0xC7, 0x59)
USER_PURPLE = RGBColor(0x7C, 0x5C, 0xD3)


def draw_icon_phone(slide, x, y, size=Inches(0.55)):
    """Stylized phone icon using a rounded rectangle + a small notch."""
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size * 0.62, size,
    )
    fill_solid(body, FLUTTER_BLUE)
    body.line.color.rgb = FLUTTER_BLUE
    body.line.width = Pt(0.5)
    # Tiny camera/notch dot at top
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + size * 0.27, y + size * 0.09,
        size * 0.08, size * 0.08,
    )
    fill_solid(dot, RGBColor(0xFF, 0xFF, 0xFF))
    set_no_line(dot)


def draw_icon_flame(slide, x, y, size=Inches(0.55)):
    """Firebase-style flame approximated with a teardrop + an inner teardrop."""
    outer = slide.shapes.add_shape(MSO_SHAPE.TEAR, x, y, size * 0.7, size)
    outer.rotation = 180  # point up
    fill_solid(outer, FIREBASE_AMBER)
    set_no_line(outer)
    inner = slide.shapes.add_shape(
        MSO_SHAPE.TEAR,
        x + size * 0.18, y + size * 0.30,
        size * 0.34, size * 0.55,
    )
    inner.rotation = 180
    fill_solid(inner, FIREBASE_YELLOW)
    set_no_line(inner)


def draw_icon_map(slide, x, y, size=Inches(0.55)):
    """Mapbox-style icon — a folded-map shape via two rectangles + a pin."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, size, size * 0.85,
    )
    fill_solid(bg, MAPBOX_BLUE)
    set_no_line(bg)
    # Crease line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + size * 0.45, y, size * 0.04, size * 0.85,
    )
    fill_solid(line, RGBColor(0xFF, 0xFF, 0xFF))
    set_no_line(line)
    # Pin dot
    pin = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + size * 0.65, y + size * 0.22,
        size * 0.18, size * 0.18,
    )
    fill_solid(pin, RGBColor(0xD2, 0x50, 0x2A))
    set_no_line(pin)


def draw_icon_pin(slide, x, y, size=Inches(0.55)):
    """Location pin — teardrop pointing down + inner circle."""
    body = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x, y, size * 0.75, size * 0.75,
    )
    fill_solid(body, GPS_GREEN)
    set_no_line(body)
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + size * 0.27, y + size * 0.27,
        size * 0.20, size * 0.20,
    )
    fill_solid(dot, RGBColor(0xFF, 0xFF, 0xFF))
    set_no_line(dot)


def draw_icon_user(slide, x, y, size=Inches(0.55)):
    """Person icon — head circle + shoulders trapezoid."""
    head = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x + size * 0.22, y, size * 0.36, size * 0.36,
    )
    fill_solid(head, USER_PURPLE)
    set_no_line(head)
    body = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x + size * 0.05, y + size * 0.40,
        size * 0.70, size * 0.45,
    )
    fill_solid(body, USER_PURPLE)
    set_no_line(body)


def draw_service_card(slide, *, x, y, w, h, title, subtitle, accent,
                      icon_drawer):
    """A white card with a colored top stripe, an icon, a title, and a
    subtitle. Used for the cloud-service boxes on the right side."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    fill_solid(card, RGBColor(0xFF, 0xFF, 0xFF))
    card.line.color.rgb = LIGHT_BORDER
    card.line.width = Pt(0.75)

    # Colored stripe at the top
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.10),
    )
    fill_solid(stripe, accent)
    set_no_line(stripe)

    # Icon
    icon_drawer(slide, x + Inches(0.20), y + Inches(0.25), Inches(0.55))

    # Title + subtitle (text starts to the right of the icon)
    text_x = x + Inches(0.95)
    text_w = w - Inches(1.05)
    tb = slide.shapes.add_textbox(text_x, y + Inches(0.22), text_w, Inches(0.85))
    tf = add_text(
        tb, title, size=12.5, bold=True, color=TITLE_GRAY, font="Calibri",
    )
    add_paragraph(tf, subtitle, size=10, color=BODY_GRAY)


# ── Slide 2: Nearby-Request architecture ──────────────────────────────────

def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Project Architecture –  Nearby Request feature")

    # Subtitle
    sub = slide.shapes.add_textbox(
        Inches(0.45), Inches(1.25), Inches(12.5), Inches(0.40),
    )
    add_text(
        sub,
        "»  How the volunteer's screen of nearby help requests is built — "
        "from the phone out to the cloud services.",
        size=13, color=BODY_GRAY,
    )

    # ── LEFT: phone container ("Mobile App / Flutter") with internal services
    phone_x, phone_y = Inches(0.45), Inches(1.85)
    phone_w, phone_h = Inches(6.30), Inches(4.20)

    # Phone outer frame
    phone = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, phone_x, phone_y, phone_w, phone_h,
    )
    fill_solid(phone, SOFT_BLUE_BG)
    phone.line.color.rgb = FLUTTER_BLUE
    phone.line.width = Pt(1.5)

    # Phone title with little phone icon
    draw_icon_phone(slide, phone_x + Inches(0.20), phone_y + Inches(0.18),
                    Inches(0.45))
    title_box = slide.shapes.add_textbox(
        phone_x + Inches(0.85), phone_y + Inches(0.12),
        phone_w - Inches(1.0), Inches(0.50),
    )
    tf = add_text(
        title_box, "Flutter Mobile App",
        size=15, bold=True, color=FLUTTER_BLUE, font="Cambria",
    )
    add_paragraph(
        tf, "the volunteer's phone",
        size=10, italic=True, color=BODY_GRAY,
    )

    # 4 internal component chips (UI screen + 3 services)
    chip_w = Inches(2.85)
    chip_h = Inches(1.30)
    gap = Inches(0.20)
    chip_top_y = phone_y + Inches(0.95)
    inner_pad = Inches(0.20)

    components = [
        # (title, subtitle, color stripe)
        ("RequestMap & Full-map\nScreens",
         "Renders the Mapbox map + scrollable\ncards of nearby requests",
         FLUTTER_BLUE),
        ("RequestService",
         "Streams the /requests collection\nfrom Firestore (live updates)",
         FIREBASE_AMBER),
        ("LocationService",
         "Geolocator → user's GPS\nposition. Used to compute distance.",
         GPS_GREEN),
        ("AuthService",
         "Looks up /users/{uid} so each\ncard shows the requester's name",
         USER_PURPLE),
    ]

    # Lay 4 chips in a 2x2 grid inside the phone
    for i, (t, s, color) in enumerate(components):
        col = i % 2
        row = i // 2
        cx = phone_x + inner_pad + col * (chip_w + gap)
        cy = chip_top_y + row * (chip_h + Inches(0.15))

        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, chip_w, chip_h,
        )
        fill_solid(chip, RGBColor(0xFF, 0xFF, 0xFF))
        chip.line.color.rgb = LIGHT_BORDER
        chip.line.width = Pt(0.75)

        # Top stripe
        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx, cy, chip_w, Inches(0.10),
        )
        fill_solid(stripe, color)
        set_no_line(stripe)

        # Title + body
        tb = slide.shapes.add_textbox(
            cx + Inches(0.15), cy + Inches(0.20),
            chip_w - Inches(0.30), chip_h - Inches(0.25),
        )
        tf = add_text(tb, t, size=11.5, bold=True, color=TITLE_GRAY)
        add_paragraph(tf, "", size=2)
        add_paragraph(tf, s, size=9.5, color=BODY_GRAY)

    # ── RIGHT: external services with shape-based icons
    svc_x = Inches(7.15)
    svc_w = Inches(5.75)
    svc_h = Inches(0.95)
    svc_gap = Inches(0.20)
    svc_y = Inches(1.85)

    services = [
        ("Cloud Firestore",
         "/requests + /users  ·  realtime stream",
         FIREBASE_AMBER, draw_icon_flame),
        ("Mapbox SDK",
         "vector map tiles  ·  custom marker pins",
         MAPBOX_NAVY, draw_icon_map),
        ("GPS  (Android Location)",
         "Geolocator plugin  ·  fine-accuracy fix",
         GPS_GREEN, draw_icon_pin),
        ("Other volunteers / requesters",
         "create requests visible to everyone nearby",
         USER_PURPLE, draw_icon_user),
    ]
    for i, (title, sub, color, drawer) in enumerate(services):
        y = svc_y + i * (svc_h + svc_gap)
        draw_service_card(
            slide,
            x=svc_x, y=y, w=svc_w, h=svc_h,
            title=title, subtitle=sub,
            accent=color, icon_drawer=drawer,
        )

    # ── Connecting arrows between phone and the 4 service cards
    for i in range(4):
        y = svc_y + i * (svc_h + svc_gap) + Inches(0.42)
        arr = slide.shapes.add_shape(
            MSO_SHAPE.LEFT_RIGHT_ARROW,
            phone_x + phone_w + Inches(0.02), y,
            svc_x - phone_x - phone_w - Inches(0.04), Inches(0.14),
        )
        fill_solid(arr, RGBColor(0x9C, 0xA3, 0xAF))
        set_no_line(arr)

    # ── Bottom: numbered flow strip
    flow_y = Inches(6.20)
    flow_steps = [
        "1.  Open screen",
        "2.  Get GPS",
        "3.  Stream  /requests",
        "4.  Resolve usernames",
        "5.  Compute distance",
        "6.  Render map + list",
    ]
    step_w = Inches(2.00)
    step_h = Inches(0.55)
    total_w = step_w * len(flow_steps) + Inches(0.10) * (len(flow_steps) - 1)
    start_x = (Inches(13.33) - total_w) / 2
    for i, label in enumerate(flow_steps):
        x = start_x + i * (step_w + Inches(0.10))
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, flow_y, step_w, step_h,
        )
        # Cycle through brand colors for visual rhythm
        colors_cycle = [FLUTTER_BLUE, GPS_GREEN, FIREBASE_AMBER,
                        USER_PURPLE, MAPBOX_BLUE, RGBColor(0xD2, 0x50, 0x2A)]
        fill_solid(chip, colors_cycle[i % len(colors_cycle)])
        set_no_line(chip)
        # Centered label text
        tf = chip.text_frame
        tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        r.font.name = "Calibri"; r.font.size = Pt(10.5); r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    add_footer(slide, 2)


def layer_box(slide, *, x, y, w, h, title, bg, chips):
    """Draws one architecture-layer band: a colored rounded background with
    a layer title on the left and 4 component cards spread across."""
    band = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h,
    )
    fill_solid(band, bg)
    band.line.color.rgb = LIGHT_BORDER
    band.line.width = Pt(0.75)

    # Layer title (vertical-left)
    title_box = slide.shapes.add_textbox(
        x + Inches(0.20), y + Inches(0.12),
        Inches(2.25), Inches(0.45),
    )
    add_text(
        title_box, title,
        size=14, bold=True, color=TITLE_GRAY, font="Cambria",
    )

    # Component chips spread across the rest of the band.
    chip_count = len(chips)
    chip_area_x = x + Inches(2.55)
    chip_area_w = w - Inches(2.75)
    chip_w = (chip_area_w - Inches(0.20) * (chip_count - 1)) / chip_count
    chip_h = h - Inches(0.65)
    chip_y = y + Inches(0.55)

    for i, (heading, body) in enumerate(chips):
        cx = chip_area_x + (chip_w + Inches(0.20)) * i
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cx, chip_y, chip_w, chip_h,
        )
        fill_solid(chip, RGBColor(0xFF, 0xFF, 0xFF))
        chip.line.color.rgb = LIGHT_BORDER
        chip.line.width = Pt(0.75)
        # Heading + body inside the chip.
        tf = chip.text_frame
        tf.margin_left = Inches(0.10)
        tf.margin_right = Inches(0.10)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = heading
        r.font.name = "Calibri"; r.font.size = Pt(12); r.font.bold = True
        r.font.color.rgb = TITLE_GRAY
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        r2.font.name = "Calibri"; r2.font.size = Pt(9.5)
        r2.font.color.rgb = BODY_GRAY


# ── Build the deck ─────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_nearby_request(prs)
    slide_architecture(prs)

    # Write a fresh copy each run so a locked PPT in PowerPoint doesn't
    # block the build. Latest version always lands at the v2 path.
    out = "helpvrywhere_two_slides_v2.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({prs.slide_width / 914400}×"
          f"{prs.slide_height / 914400} in, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
